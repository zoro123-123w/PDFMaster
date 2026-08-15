"""Real conversion services: PDF <-> Word, PDF -> Excel, PDF -> PPT, HTML/TXT -> PDF.

Each function takes a saved file path plus the form cleaned_data and returns the
path of a generated output file inside MEDIA_ROOT/temp.
"""
import os
import tempfile

from .utils import create_output_path

# ---------------------------------------------------------------------------
# PDF -> Word
# ---------------------------------------------------------------------------

def _extract_pdf_images(fitz_page, temp_path):
    """Save every image found on a page to temp files; return their paths."""
    saved = []
    for img in fitz_page.get_images(full=True):
        try:
            xref = img[0]
            pix = fitz_page.parent.extract_image(xref)
            if pix is None:
                continue
            img_bytes = pix['image']
            ext = pix.get('ext', 'png')
            path = os.path.join(temp_path, f'{xref}_{len(saved)}.{ext}')
            with open(path, 'wb') as fh:
                fh.write(img_bytes)
            saved.append(path)
        except Exception:
            continue
    return saved


def pdf_to_word(file_path, data=None):
    """Convert a PDF to a .docx file. Text is written as paragraphs and images
    embedded so the document is genuinely editable in Word."""
    import fitz
    from docx import Document
    from docx.shared import Inches

    output_path = create_output_path('.docx', 'pdf2word_')
    pdf = fitz.open(file_path)
    doc = Document()
    tmp_img_dir = tempfile.mkdtemp(prefix='pdf2word_')
    try:
        for page in pdf:
            text = page.get_text()
            if text and text.strip():
                lines = [l.strip() for l in text.splitlines() if l.strip()]
                for line in lines:
                    doc.add_paragraph(line)
            images = _extract_pdf_images(page, tmp_img_dir)
            for img_path in images[:10]:  # cap per page to avoid huge files
                try:
                    doc.add_picture(img_path, width=Inches(6))
                    doc.add_paragraph('')
                except Exception:
                    continue
            doc.add_page_break()
        doc.save(output_path)
    finally:
        pdf.close()
        for root, _dirs, files in os.walk(tmp_img_dir):
            for f in files:
                try:
                    os.remove(os.path.join(root, f))
                except OSError:
                    pass
        try:
            os.rmdir(tmp_img_dir)
        except OSError:
            pass
    return output_path


# ---------------------------------------------------------------------------
# Word -> PDF
# ---------------------------------------------------------------------------

def word_to_pdf(file_path, data=None):
    """Convert a .docx file to a PDF using reportlab. Real text conversion:
    paragraphs, headings and tables are rendered into the PDF."""
    from docx import Document
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.units import inch
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table

    output_path = create_output_path('.pdf', 'word2pdf_')
    doc = Document(file_path)
    styles = getSampleStyleSheet()
    body = ParagraphStyle('BodyTextEx', parent=styles['BodyText'], spaceAfter=8)

    h_styles = {}
    for level in range(1, 5):
        h_styles[level] = ParagraphStyle(
            f'HeadingEx{level}', parent=styles['Heading%d' % level],
            spaceBefore=12, spaceAfter=6)

    story = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        name = (para.style.name or '').lower()
        if name.startswith('heading'):
            try:
                level = int(name.split()[-1])
            except (ValueError, IndexError):
                level = 1
            if '1' in name and level > 1:
                level = 1
            style = h_styles.get(min(level, 4), h_styles[1])
            story.append(Paragraph(_xml_escape(text), style))
        else:
            story.append(Paragraph(_xml_escape(text), body))

    for table in doc.tables:
        rows_data = []
        for row in table.rows:
            cells = [Paragraph(_xml_escape((c.text or '').strip()[:1000]), body)
                     for c in row.cells]
            rows_data.append(cells)
        if rows_data:
            story.append(Table(rows_data, repeatRows=1))
            story.append(Spacer(1, 12))

    pdf_doc = SimpleDocTemplate(output_path, pagesize=letter,
                                leftMargin=inch, rightMargin=inch,
                                topMargin=inch, bottomMargin=inch)
    pdf_doc.build(story)
    return output_path


def _xml_escape(text):
    return (text.replace('&', '&amp;').replace('<', '&lt;')
                .replace('>', '&gt;'))


# ---------------------------------------------------------------------------
# PDF -> Excel
# ---------------------------------------------------------------------------

def pdf_to_excel(file_path, data=None):
    """Convert a PDF to an .xlsx workbook. Tables extracted with pdfplumber go
    into worksheets; free text is appended to a 'PDF Content' sheet."""
    import pdfplumber
    from openpyxl import Workbook

    output_path = create_output_path('.xlsx', 'pdf2excel_')
    wb = Workbook()

    with pdfplumber.open(file_path) as pdf:
        text_sheet = wb.active
        text_sheet.title = 'PDF Content'
        text_row = 1

        for page_index, page in enumerate(pdf.pages, start=1):
            tables = []
            try:
                tables = page.extract_tables() or []
            except Exception:
                tables = []
            for t_idx, table in enumerate(tables):
                sheet = wb.create_sheet(title=f'Page{page_index}_Table{t_idx + 1}')
                max_len = max((len(r or []) for r in table), default=0)
                for row in table:
                    cols = [(c or '') for c in (row or [])]
                    while len(cols) < max_len:
                        cols.append('')
                    sheet.append(cols)

            page_text = page.extract_text()
            if page_text:
                if text_row > 1:
                    text_sheet.cell(row=text_row, column=1, value=f'--- Page {page_index} ---')
                    text_row += 1
                for line in page_text.splitlines():
                    text_sheet.cell(row=text_row, column=1, value=line)
                    text_row += 1

    wb.save(output_path)
    return output_path
# ---------------------------------------------------------------------------
# PDF -> PowerPoint
# ---------------------------------------------------------------------------

def pdf_to_ppt(file_path, data=None):
    """Convert each PDF page to a slide containing a rendered page image plus
    its extracted text, producing a real .pptx file."""
    import fitz
    from pptx import Presentation
    from pptx.util import Inches

    output_path = create_output_path('.pptx', 'pdf2ppt_')
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_w_in = 13.333
    slide_h_in = 7.5

    pdf = fitz.open(file_path)
    tmp_img_dir = tempfile.mkdtemp(prefix='pdf2ppt_')
    try:
        for page_index, page in enumerate(pdf):
            pix = page.get_pixmap(dpi=80)
            img_path = os.path.join(tmp_img_dir, f'page_{page_index}.png')
            pix.save(img_path)

            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
            aspect = pix.width / pix.height if pix.height else 1.5
            img_w = min(7.8, slide_h_in * 0.92 * aspect)
            img_h = img_w / aspect
            left = Inches(0.2)
            top = Inches((slide_h_in - img_h) / 2)
            try:
                slide.shapes.add_picture(img_path, left, top,
                                         width=Inches(img_w), height=Inches(img_h))
            except Exception:
                pass

            text = page.get_text().strip()
            if text:
                text_box = slide.shapes.add_textbox(Inches(8.3), Inches(0.3),
                                                    Inches(4.8), Inches(6.8))
                tf = text_box.text_frame
                tf.word_wrap = True
                tf.text = text[:1800]
    finally:
        pdf.close()
        for root, _dirs, files in os.walk(tmp_img_dir):
            for f in files:
                try:
                    os.remove(os.path.join(root, f))
                except OSError:
                    pass
        try:
            os.rmdir(tmp_img_dir)
        except OSError:
            pass

    prs.save(output_path)
    return output_path


# ---------------------------------------------------------------------------
# HTML -> PDF
# ---------------------------------------------------------------------------

def html_to_pdf(file_path, data=None):
    """Convert an HTML file to a PDF using xhtml2pdf (pisa)."""
    from xhtml2pdf import pisa

    output_path = create_output_path('.pdf', 'html2pdf_')
    with open(file_path, 'rb') as src:
        html_data = src.read()
    with open(output_path, 'wb') as dest:
        status = pisa.CreatePDF(html_data, dest=dest,
                                encoding='utf-8',
                                base_url=os.path.dirname(file_path))
    if status.err:
        raise ValueError('Could not render HTML to PDF - invalid or unsupported HTML content.')
    return output_path


# ---------------------------------------------------------------------------
# TXT -> PDF
# ---------------------------------------------------------------------------

def txt_to_pdf(file_path, data=None):
    """Convert a plain text file to a paginated PDF using reportlab."""
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Preformatted

    output_path = create_output_path('.pdf', 'txt2pdf_')

    with open(file_path, 'r', encoding='utf-8', errors='replace') as fh:
        raw = fh.read()

    styles = getSampleStyleSheet()
    code_style = ParagraphStyle('TxtEx', parent=styles['BodyText'],
                                fontName='Courier', fontSize=10,
                                leading=13, spaceAfter=6)
    story = []
    for para in raw.split('\n\n'):
        escaped = para.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
        if escaped.strip():
            story.append(Preformatted(escaped.strip(), code_style))

    doc = SimpleDocTemplate(output_path, pagesize=letter,
                            leftMargin=inch, rightMargin=inch,
                            topMargin=inch, bottomMargin=inch)
    doc.build(story)
    return output_path


# ---------------------------------------------------------------------------
# Excel -> PDF
# ---------------------------------------------------------------------------

def excel_to_pdf(file_path, data=None):
    """Convert each worksheet of an .xlsx workbook to PDF pages using
    openpyxl for data extraction and reportlab for rendering."""
    from openpyxl import load_workbook
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.units import inch
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib import colors
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                    Table, TableStyle)

    output_path = create_output_path('.pdf', 'excel2pdf_')
    wb = load_workbook(file_path, data_only=True, read_only=True)

    styles = getSampleStyleSheet()
    body = ParagraphStyle('ExcelBody', parent=styles['BodyText'], fontSize=8, spaceAfter=4)
    heading = ParagraphStyle('ExcelHeading', parent=styles['Heading2'],
                             spaceBefore=12, spaceAfter=6)

    story = []

    for ws in wb.worksheets:
        title = ws.title or 'Sheet'
        story.append(Paragraph(title, heading))
        rows_data = []
        for row in ws.iter_rows(values_only=True):
            cells = [(str(c) if c is not None else '') for c in row]
            rows_data.append(cells)
        if rows_data:
            max_cols = max(len(r) for r in rows_data)
            for r in rows_data:
                while len(r) < max_cols:
                    r.append('')
            table = Table(rows_data, repeatRows=1)
            table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#e2e8f0')),
                ('FONTSIZE', (0, 0), (-1, -1), 8),
                ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                ('GRID', (0, 0), (-1, -1), 0.25, colors.HexColor('#cbd5e1')),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 2),
                ('TOPPADDING', (0, 0), (-1, -1), 2),
            ]))
            story.append(table)
        story.append(Spacer(1, 12))

    wb.close()
    pdf_doc = SimpleDocTemplate(output_path, pagesize=letter,
                                leftMargin=inch, rightMargin=inch,
                                topMargin=inch, bottomMargin=inch)
    pdf_doc.build(story)
    return output_path


# ---------------------------------------------------------------------------
# PowerPoint -> PDF
# ---------------------------------------------------------------------------

def pptx_to_pdf(file_path, data=None):
    """Convert each PowerPoint slide to a PDF page using reportlab.
    Slide text content is rendered as paragraphs with title styling."""
    from pptx import Presentation
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.units import inch
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak

    output_path = create_output_path('.pdf', 'pptx2pdf_')
    prs = Presentation(file_path)

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle('PPTTitle', parent=styles['Heading1'],
                                 spaceAfter=12)
    body_style = ParagraphStyle('PPTBody', parent=styles['BodyText'],
                                spaceAfter=6)

    story = []
    for slide in prs.slides:
        title = ''
        if slide.shapes.title and slide.shapes.title.text:
            title = slide.shapes.title.text
        if title:
            story.append(Paragraph(_xml_escape(title), title_style))
        body_texts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text and shape != slide.shapes.title:
                body_texts.append(_xml_escape(shape.text.strip()))
        for bt in body_texts:
            for line in bt.splitlines():
                line = line.strip()
                if line:
                    story.append(Paragraph(line, body_style))
        story.append(PageBreak())

    pdf_doc = SimpleDocTemplate(output_path, pagesize=letter,
                                leftMargin=inch, rightMargin=inch,
                                topMargin=inch, bottomMargin=inch)
    pdf_doc.build(story)
    return output_path


# ---------------------------------------------------------------------------
# Scan / image -> Searchable PDF (rasterised single page)
# ---------------------------------------------------------------------------

def image_to_pdf(file_paths, data=None):
    """Combine one or more image files into a single PDF (one image per page).
    Uses Pillow for image loading so all common formats are supported."""
    from PIL import Image
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter

    if not file_paths:
        raise ValueError('No images provided.')

    output_path = create_output_path('.pdf', 'images2pdf_')
    c = canvas.Canvas(output_path, pagesize=letter)
    page_width, page_height = letter

    for img_path in file_paths:
        img = Image.open(img_path)
        if img.mode in ('RGBA', 'LA', 'P'):
            img = img.convert('RGB')
        width, height = img.size
        aspect = height / float(width) if width else 1
        if width > height:
            new_width = page_width
            new_height = page_width * aspect
        else:
            new_height = page_height
            new_width = page_height / aspect if aspect else page_height
        if new_width > page_width:
            new_width = page_width
            new_height = page_width * aspect
        if new_height > page_height:
            new_height = page_height
            new_width = page_height / aspect if aspect else page_height
        c.setPageSize((page_width, page_height))
        c.drawImage(img_path, (page_width - new_width) / 2,
                    (page_height - new_height) / 2,
                    width=new_width, height=new_height)
        c.showPage()
    c.save()
    return output_path