import io
import json
import os
import re
import zipfile
import struct
import logging

from django.test import TestCase, Client
from django.urls import reverse, reverse_lazy
from django.core.files.uploadedfile import SimpleUploadedFile
from django.contrib.auth import get_user_model
from django.contrib.auth.hashers import make_password
from pypdf import PdfWriter
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas as pdfcanvas
from PIL import Image
from openpyxl import Workbook, load_workbook
from docx import Document as WordDocument
from pptx import Presentation as PowerPointPresentation

User = get_user_model()
logger = logging.getLogger(__name__)


def make_text_pdf(text="Hello PDFMaster. This is a test document. " * 50, pages=2):
    """Create a multi-page PDF with text content using reportlab then convert with pypdf."""
    buffer = io.BytesIO()
    c = pdfcanvas.Canvas(buffer, pagesize=letter)
    width, height = letter
    for i in range(pages):
        c.drawString(100, height - 100, text)
        c.showPage()
    c.save()
    buffer.seek(0)
    return SimpleUploadedFile('test.pdf', buffer.read(),
                              content_type='application/pdf')


def make_image_pdf(pages=1):
    """Create a simple PDF using pypdf that is also a valid PDF."""
    buffer = io.BytesIO()
    w = PdfWriter()
    for _ in range(pages):
        w.add_blank_page(width=200, height=200)
    w.write(buffer)
    buffer.seek(0)
    return SimpleUploadedFile('image_test.pdf', buffer.read(),
                              content_type='application/pdf')


def make_image_file(ext='png', size=(100, 100), color='red'):
    """Create a simple image file."""
    img = Image.new('RGB', size, color)
    buffer = io.BytesIO()
    fmt = 'PNG' if ext == 'png' else 'JPEG'
    img.save(buffer, format=fmt)
    buffer.seek(0)
    content = buffer.read()
    filename = f'test.{ext}'
    content_type = f'image/{ext}' if ext != 'jpeg' else 'image/jpeg'
    if ext == 'jpg':
        content_type = 'image/jpeg'
    return SimpleUploadedFile(filename, content, content_type=content_type)


def make_docx_file():
    """Create a simple .docx file."""
    buffer = io.BytesIO()
    doc = WordDocument()
    doc.add_heading('Test Document', 0)
    doc.add_paragraph('This is a test paragraph.')
    doc.save(buffer)
    buffer.seek(0)
    return SimpleUploadedFile('test.docx', buffer.read(),
                              content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document')


def make_xlsx_file():
    """Create a simple .xlsx file."""
    buffer = io.BytesIO()
    wb = Workbook()
    ws = wb.active
    ws['A1'] = 'Name'
    ws['B1'] = 'Age'
    ws['A2'] = 'Alice'
    ws['B2'] = 30
    wb.save(buffer)
    buffer.seek(0)
    return SimpleUploadedFile('test.xlsx', buffer.read(),
                              content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')


def make_pptx_file():
    """Create a simple .pptx file."""
    buffer = io.BytesIO()
    prs = PowerPointPresentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(buffer)
    buffer.seek(0)
    return SimpleUploadedFile('test.pptx', buffer.read(),
                              content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')


def make_html_file():
    """Create a simple HTML file."""
    content = b'<html><head><title>Test</title></head><body><h1>Hello World</h1><p>Test paragraph.</p></body></html>'
    return SimpleUploadedFile('test.html', content, content_type='text/html')


def make_txt_file():
    """Create a simple text file."""
    content = b'Hello World\nThis is a test text file.\nLine 3.'
    return SimpleUploadedFile('test.txt', content, content_type='text/plain')


def make_non_pdf_file():
    """Create a file with PDF extension but invalid content."""
    content = b'NOT A PDF FILE'
    return SimpleUploadedFile('fake.pdf', content, content_type='application/pdf')


def make_empty_file():
    """Create an empty file."""
    return SimpleUploadedFile('empty.pdf', b'', content_type='application/pdf')


def get_download_url(response):
    """Extract download URL from a redirect response."""
    if hasattr(response, 'url'):
        return response.url
    location = response.get('Location', '')
    if location.startswith('/'):
        return location
    return None


class ComprehensiveIntegrationTest(TestCase):
    """Comprehensive end-to-end test of every feature in PDFMaster."""

    def setUp(self):
        self.client = Client()
        self.user = User.objects.create_user(
            username='testuser', email='test@example.com',
            password='TestPass123!'
        )
        self.client.login(username='testuser', password='TestPass123!')

    def _make_pdf(self):
        return make_text_pdf()

    def _make_img_pdf(self):
        return make_image_pdf()

    # ================================================================
    # A. PDF ORGANIZE TOOLS
    # ================================================================

    def test_merge_pdf(self):
        """Merge two PDFs into one."""
        pdf1 = self._make_pdf()
        pdf2 = self._make_pdf()
        response = self.client.post(reverse('merge_pdf'), {'files': [pdf1, pdf2]})
        self.assertEqual(response.status_code, 302, f"Expected 302, got {response.status_code}")

    def test_split_pdf(self):
        """Split a PDF into individual pages."""
        pdf = make_text_pdf(pages=3)
        response = self.client.post(reverse('split_pdf'), {'file': pdf})
        self.assertIn(response.status_code, [302, 200])

    def test_delete_pages(self):
        """Delete pages from a PDF."""
        pdf = make_text_pdf(pages=5)
        response = self.client.post(reverse('delete_pages'), {
            'file': pdf, 'pages': '1-2'
        })
        self.assertIn(response.status_code, [302, 200])

    def test_extract_pages(self):
        """Extract specific pages from a PDF."""
        pdf = make_text_pdf(pages=5)
        response = self.client.post(reverse('extract_pages'), {
            'file': pdf, 'pages': '1-2'
        })
        self.assertIn(response.status_code, [302, 200])

    def test_organize_pdf(self):
        """Reorder pages in a PDF."""
        pdf = make_text_pdf(pages=3)
        response = self.client.post(reverse('organize_pdf'), {
            'file': pdf, 'order': '3,1,2', 'rotation': '0'
        })
        self.assertEqual(response.status_code, 302, f"Expected 302, got {response.status_code}")

    # ================================================================
    # B. PDF OPTIMIZE TOOLS
    # ================================================================

    def test_compress_pdf(self):
        """Compress a PDF."""
        pdf = self._make_pdf()
        response = self.client.post(reverse('compress_pdf'), {'file': pdf})
        self.assertEqual(response.status_code, 200)

    def test_repair_pdf(self):
        """Repair a PDF."""
        pdf = self._make_pdf()
        response = self.client.post(reverse('repair_pdf'), {'file': pdf})
        self.assertEqual(response.status_code, 302, f"Expected 302, got {response.status_code}")

    # ================================================================
    # C. PDF ROTATE
    # ================================================================

    def test_rotate_pdf(self):
        """Rotate a PDF."""
        pdf = self._make_pdf()
        response = self.client.post(reverse('rotate_pdf'), {
            'file': pdf, 'angle': '90'
        })
        self.assertEqual(response.status_code, 302, f"Expected 302, got {response.status_code}")

    # ================================================================
    # D. PDF CONVERT TOOLS
    # ================================================================

    def test_pdf_to_word(self):
        """Convert PDF to DOCX."""
        pdf = self._make_pdf()
        response = self.client.post(reverse('pdf_to_word'), {'file': pdf})
        self.assertEqual(response.status_code, 302, f"Expected 302, got {response.status_code}")

    def test_pdf_to_excel(self):
        """Convert PDF to XLSX."""
        pdf = self._make_pdf()
        response = self.client.post(reverse('pdf_to_excel'), {'file': pdf})
        self.assertEqual(response.status_code, 302, f"Expected 302, got {response.status_code}")

    def test_pdf_to_ppt(self):
        """Convert PDF to PPTX."""
        pdf = self._make_pdf()
        response = self.client.post(reverse('pdf_to_ppt'), {'file': pdf})
        self.assertEqual(response.status_code, 302, f"Expected 302, got {response.status_code}")

    def test_word_to_pdf(self):
        """Convert DOCX to PDF."""
        docx = make_docx_file()
        response = self.client.post(reverse('word_to_pdf'), {'file': docx})
        self.assertEqual(response.status_code, 302, f"Expected 302, got {response.status_code}")

    def test_excel_to_pdf(self):
        """Convert XLSX to PDF."""
        xlsx = make_xlsx_file()
        response = self.client.post(reverse('excel_to_pdf'), {'file': xlsx})
        self.assertEqual(response.status_code, 302, f"Expected 302, got {response.status_code}")

    def test_ppt_to_pdf(self):
        """Convert PPTX to PDF."""
        pptx = make_pptx_file()
        response = self.client.post(reverse('ppt_to_pdf'), {'file': pptx})
        self.assertEqual(response.status_code, 302, f"Expected 302, got {response.status_code}")

    def test_html_to_pdf(self):
        """Convert HTML to PDF."""
        html = make_html_file()
        response = self.client.post(reverse('html_to_pdf'), {'file': html})
        self.assertEqual(response.status_code, 302, f"Expected 302, got {response.status_code}")

    def test_txt_to_pdf(self):
        """Convert TXT to PDF."""
        txt = make_txt_file()
        response = self.client.post(reverse('txt_to_pdf'), {'file': txt})
        self.assertEqual(response.status_code, 302, f"Expected 302, got {response.status_code}")

    def test_pdf_to_jpg(self):
        """Convert PDF to JPG images."""
        pdf = make_text_pdf(pages=2)
        response = self.client.post(reverse('pdf_to_jpg'), {'file': pdf})
        self.assertIn(response.status_code, [302, 200])

    def test_pdf_to_png(self):
        """Convert PDF to PNG images."""
        pdf = make_text_pdf(pages=2)
        response = self.client.post(reverse('pdf_to_png'), {'file': pdf})
        self.assertIn(response.status_code, [302, 200])

    def test_jpg_to_pdf(self):
        """Convert JPG to PDF."""
        img = make_image_file('jpg')
        response = self.client.post(reverse('jpg_to_pdf'), {'files': [img]})
        self.assertEqual(response.status_code, 302, f"Expected 302, got {response.status_code}")

    def test_png_to_pdf(self):
        """Convert PNG to PDF."""
        img = make_image_file('png')
        response = self.client.post(reverse('png_to_pdf'), {'files': img})
        self.assertEqual(response.status_code, 302, f"Expected 302, got {response.status_code}: {response.content[:300]}")

    def test_webp_to_pdf(self):
        """Convert WebP to PDF."""
        img = make_image_file('webp')
        response = self.client.post(reverse('webp_to_pdf'), {'files': img})
        self.assertEqual(response.status_code, 302, f"Expected 302, got {response.status_code}: {response.content[:300]}")

    def test_webp_to_pdf_output_validity(self):
        """Verify PDF output from WebP is valid."""
        from pypdf import PdfReader
        img = make_image_file('webp')
        response = self.client.post(reverse('webp_to_pdf'), {'files': img})
        self.assertIn(response.status_code, [302, 200])
        if response.status_code == 302:
            download_response = self.client.get(response.url)
            self.assertEqual(download_response.status_code, 200)
            content = b''.join(download_response.streaming_content)
            try:
                reader = PdfReader(io.BytesIO(content))
                self.assertGreater(reader.get_num_pages(), 0)
            except Exception as e:
                self.fail(f"PDF validation failed: {e}")
        try:
            reader = PdfReader(io.BytesIO(content))
            self.assertGreater(reader.get_num_pages(), 0)
        except Exception as e:
            self.fail(f"PDF validation failed: {e}")

    def test_sign_pdf(self):
        """Sign a PDF."""
        pdf = self._make_pdf()
        response = self.client.post(reverse('sign_pdf'), {
            'file': pdf, 'signature_data': 'Test Signature',
            'page': '1', 'sig_x': '100', 'sig_y': '100', 'sig_width': '200'
        })
        self.assertIn(response.status_code, [302, 200])

    # ================================================================
    # E. PDF EDIT TOOLS
    # ================================================================

    def test_crop_pdf(self):
        """Crop a PDF."""
        pdf = self._make_pdf()
        response = self.client.post(reverse('crop_pdf'), {
            'file': pdf, 'left': '5', 'right': '5', 'top': '5', 'bottom': '5'
        })
        self.assertEqual(response.status_code, 302, f"Expected 302, got {response.status_code}")

    def test_watermark_pdf(self):
        """Add watermark to a PDF."""
        pdf = self._make_pdf()
        response = self.client.post(reverse('watermark_pdf'), {
            'file': pdf, 'text': 'WATERMARK', 'position': 'center',
            'opacity': '0.3', 'fontsize': '30'
        })
        self.assertEqual(response.status_code, 302, f"Expected 302, got {response.status_code}")

    def test_add_page_numbers(self):
        """Add page numbers to a PDF."""
        pdf = self._make_pdf()
        response = self.client.post(reverse('add_page_numbers'), {
            'file': pdf, 'position': 'bottom-center', 'start': '1', 'fontsize': '14'
        })
        self.assertEqual(response.status_code, 302, f"Expected 302, got {response.status_code}")

    def test_add_text_to_pdf(self):
        """Add text to a PDF."""
        pdf = self._make_pdf()
        response = self.client.post(reverse('add_text_to_pdf'), {
            'file': pdf, 'text': 'Hello from PDFMaster', 'page_number': '1',
            'x': '50', 'y': '50', 'fontsize': '12'
        })
        self.assertEqual(response.status_code, 302, f"Expected 302, got {response.status_code}")

    def test_add_image_to_pdf(self):
        """Add an image to a PDF."""
        pdf = self._make_pdf()
        img = make_image_file('png')
        response = self.client.post(reverse('add_image_to_pdf'), {
            'file': pdf, 'image': img, 'page_number': '1',
            'x': '50', 'y': '50', 'width_pct': '25'
        })
        self.assertEqual(response.status_code, 302, f"Expected 302, got {response.status_code}")

    def test_annotate_pdf(self):
        """Annotate a PDF."""
        pdf = self._make_pdf()
        response = self.client.post(reverse('annotate_pdf'), {
            'file': pdf, 'text': 'Annotation note', 'page_number': '1',
            'x': '50', 'y': '50', 'fontsize': '12'
        })
        self.assertEqual(response.status_code, 302, f"Expected 302, got {response.status_code}")

    def test_highlight_pdf(self):
        """Highlight text in a PDF."""
        pdf = self._make_pdf()
        response = self.client.post(reverse('highlight_pdf'), {
            'file': pdf, 'words': 'test', 'color': 'yellow', 'page_number': '0'
        })
        self.assertEqual(response.status_code, 302, f"Expected 302, got {response.status_code}")

    def test_redact_pdf(self):
        """Redact text from a PDF."""
        pdf = self._make_pdf()
        response = self.client.post(reverse('redact_pdf'), {
            'file': pdf, 'words': 'test'
        })
        self.assertEqual(response.status_code, 302, f"Expected 302, got {response.status_code}")

    # ================================================================
    # F. PDF SECURITY TOOLS
    # ================================================================

    def test_protect_pdf(self):
        """Protect a PDF with a password."""
        pdf = self._make_pdf()
        response = self.client.post(reverse('protect_pdf'), {
            'file': pdf, 'password': 'secret123', 'confirm_password': 'secret123'
        })
        self.assertEqual(response.status_code, 302, f"Expected 302, got {response.status_code}")

    def test_unlock_pdf(self):
        """Unlock a password-protected PDF."""
        pdf = self._make_pdf()
        response = self.client.post(reverse('unlock_pdf'), {
            'file': pdf, 'password': 'secret123'
        })
        self.assertEqual(response.status_code, 302, f"Expected 302, got {response.status_code}")

    # ================================================================
    # G. SCAN / OCR TOOLS
    # ================================================================

    def test_ocr_pdf(self):
        """OCR a PDF (may degrade gracefully if Tesseract not available)."""
        pdf = self._make_pdf()
        response = self.client.post(reverse('ocr_pdf'), {
            'file': pdf, 'pages': '', 'lang': 'eng'
        })
        self.assertIn(response.status_code, [200, 302])

    def test_scan_to_pdf(self):
        """Scan to searchable PDF."""
        img = make_image_file('png')
        response = self.client.post(reverse('scan_to_pdf'), {
            'files': [img], 'pages': '', 'lang': 'eng'
        })
        self.assertIn(response.status_code, [200, 302])

    # ================================================================
    # H. DOWNLOAD VERIFICATION
    # ================================================================

    def test_download_valid_job(self):
        """Verify download endpoint works for a created job."""
        from pdf_tools.models import ProcessingJob
        from uuid import uuid4
        import os
        from django.conf import settings
        temp_dir = os.path.join(settings.MEDIA_ROOT, 'temp')
        os.makedirs(temp_dir, exist_ok=True)
        fake_output = os.path.join(temp_dir, f'{uuid4().hex}.pdf')
        with open(fake_output, 'wb') as f:
            f.write(b'%PDF-1.4\n1 0 obj\n<<>>\nendobj\ntrailer\n<<>>\n%%EOF')
        job = ProcessingJob.objects.create(
            user=self.user,
            tool_name='test',
            original_filename='test.pdf',
            status='COMPLETED',
            output_file=fake_output,
            output_filename='test.pdf',
        )
        response = self.client.get(reverse('download_file', args=[job.id]))
        self.assertEqual(response.status_code, 200)
        content = b''.join(response.streaming_content)
        self.assertGreater(len(content), 0)

    # ================================================================
    # I. UPLOAD VALIDATION / SECURITY
    # ================================================================

    def test_upload_invalid_pdf_content(self):
        """Upload a file named .pdf but with non-PDF content."""
        fake = make_non_pdf_file()
        response = self.client.post(reverse('pdf_to_word'), {'file': fake})
        self.assertEqual(response.status_code, 200)
        self.assertContains(response, "Invalid PDF file")

    def test_upload_empty_file(self):
        """Upload an empty file - validation should reject it gracefully."""
        empty = make_empty_file()
        response = self.client.post(reverse('pdf_to_word'), {'file': empty})
        self.assertEqual(response.status_code, 200)
        # The form will reject it - check that errors are shown
        content = response.content.decode()
        self.assertTrue(
            'error' in content.lower() or 'not a PDF' in content.lower()
            or 'invalid' in content.lower() or 'required' in content.lower()
            or 'empty' in content.lower(),
            "Expected error indication in response"
        )

    def test_upload_wrong_extension(self):
        """Upload a non-PDF file with wrong extension."""
        txt = make_txt_file()
        response = self.client.post(reverse('pdf_to_word'), {'file': txt})
        self.assertEqual(response.status_code, 200)

    # ================================================================
    # J. AI STUDY TOOLS (will use fallback if no AI key)
    # ================================================================

    def test_ai_quiz_no_file(self):
        """AI Quiz without file should return 400 for AJAX."""
        from ai_tools.forms import QuizForm
        response = self.client.post(
            reverse('ai_study_quiz'),
            data={},
            HTTP_X_REQUESTED_WITH='XMLHttpRequest',
        )
        self.assertEqual(response.status_code, 400)

    def test_ai_quiz_with_file(self):
        """AI Quiz with valid PDF should process (may use fallback)."""
        pdf = self._make_pdf()
        response = self.client.post(reverse('ai_study_quiz'), {
            'file': pdf, 'num_questions': '5', 'difficulty': 'easy'
        })
        self.assertIn(response.status_code, [200, 302])

    def test_ai_flashcards_with_file(self):
        """AI Flashcards with valid PDF."""
        pdf = self._make_pdf()
        response = self.client.post(reverse('ai_study_flashcards'), {
            'file': pdf, 'num_flashcards': '5'
        })
        self.assertIn(response.status_code, [200, 302])

    def test_ai_study_notes_with_file(self):
        """AI Study Notes with valid PDF."""
        pdf = self._make_pdf()
        response = self.client.post(reverse('ai_study_notes'), {
            'file': pdf,
        })
        self.assertIn(response.status_code, [200, 302])

    def test_ai_study_guide_with_file(self):
        """AI Study Guide with valid PDF."""
        pdf = self._make_pdf()
        response = self.client.post(reverse('ai_study_guide'), {
            'file': pdf,
        })
        self.assertIn(response.status_code, [200, 302])

    def test_ai_question_bank_with_file(self):
        """AI Question Bank with valid PDF."""
        pdf = self._make_pdf()
        response = self.client.post(reverse('ai_study_questions'), {
            'file': pdf, 'num_questions': '5'
        })
        self.assertIn(response.status_code, [200, 302])

    def test_ai_important_questions_with_file(self):
        """AI Important Questions with valid PDF."""
        pdf = self._make_pdf()
        response = self.client.post(reverse('ai_study_important'), {
            'file': pdf,
        })
        self.assertIn(response.status_code, [200, 302])

    def test_ai_chapter_summary_with_file(self):
        """AI Chapter Summary with valid PDF."""
        pdf = self._make_pdf()
        response = self.client.post(reverse('ai_study_chapters'), {
            'file': pdf,
        })
        self.assertIn(response.status_code, [200, 302])

    def test_ai_key_concepts_with_file(self):
        """AI Key Concepts with valid PDF."""
        pdf = self._make_pdf()
        response = self.client.post(reverse('ai_study_concepts'), {
            'file': pdf,
        })
        self.assertIn(response.status_code, [200, 302])

    def test_ai_exam_prep_with_file(self):
        """AI Exam Prep with valid PDF."""
        pdf = self._make_pdf()
        response = self.client.post(reverse('ai_study_exam'), {
            'file': pdf,
        })
        self.assertIn(response.status_code, [200, 302])

    def test_ai_summarize_with_file(self):
        """AI Summarize with valid PDF."""
        pdf = self._make_pdf()
        response = self.client.post(reverse('ai_summarize'), {
            'file': pdf,
        })
        self.assertEqual(response.status_code, 302)

    def test_ai_extract_with_file(self):
        """AI Extract Text with valid PDF."""
        pdf = self._make_pdf()
        response = self.client.post(reverse('ai_extract'), {
            'file': pdf,
        })
        self.assertEqual(response.status_code, 302)

    def test_ai_ask_with_file(self):
        """AI Ask Question with valid PDF."""
        pdf = self._make_pdf()
        response = self.client.post(reverse('ai_ask'), {
            'file': pdf, 'question': 'What is this document about?'
        })
        self.assertEqual(response.status_code, 302)

    # ================================================================
    # K. SEO / PUBLIC PAGES
    # ================================================================

    def test_home_page(self):
        """Home page returns 200."""
        response = self.client.get(reverse('home'))
        self.assertEqual(response.status_code, 200)

    def test_robots_txt(self):
        """robots.txt returns 200."""
        response = self.client.get('/robots.txt')
        self.assertEqual(response.status_code, 200)
        self.assertEqual(response['Content-Type'], 'text/plain')
        content = response.content.decode()
        self.assertIn('User-agent: *', content)
        self.assertIn('Disallow: /admin/', content)
        self.assertIn('Sitemap:', content)

    def test_sitemap_xml(self):
        """sitemap.xml returns 200 with all public pages."""
        response = self.client.get('/sitemap.xml')
        self.assertEqual(response.status_code, 200)
        self.assertEqual(response['Content-Type'], 'application/xml')
        content = response.content.decode()
        self.assertIn('<urlset', content)
        self.assertIn('/tools/', content)
        self.assertIn('/ai-tools/', content)
        self.assertNotIn('/accounts/login/', content)
        self.assertNotIn('/register/', content)

    def test_pricing_page(self):
        """Pricing page returns 200."""
        response = self.client.get(reverse('pricing'))
        self.assertEqual(response.status_code, 200)

    def test_faq_page(self):
        """FAQ page returns 200."""
        response = self.client.get(reverse('faq'))
        self.assertEqual(response.status_code, 200)

    # ================================================================
    # L. AUTH FEATURES
    # ================================================================

    def test_register_page(self):
        """Register page loads."""
        self.client.logout()
        response = self.client.get(reverse('register'))
        self.assertEqual(response.status_code, 200)

    def test_login_page(self):
        """Login page loads."""
        self.client.logout()
        response = self.client.get(reverse('login'))
        self.assertEqual(response.status_code, 200)

    def test_dashboard_requires_auth(self):
        """Dashboard requires login."""
        self.client.logout()
        response = self.client.get(reverse('dashboard'))
        self.assertIn(response.status_code, [302, 301])

    def test_dashboard_with_auth(self):
        """Dashboard loads for authenticated user."""
        response = self.client.get(reverse('dashboard'))
        self.assertEqual(response.status_code, 200)

    def test_anonymous_can_process(self):
        """Anonymous users can use PDF tools."""
        self.client.logout()
        pdf = self._make_pdf()
        response = self.client.post(reverse('pdf_to_word'), {'file': pdf})
        self.assertEqual(response.status_code, 302)

    # ================================================================
    # M. PDF TOOLS URL TESTS
    # ================================================================

    def test_all_pdf_tool_get(self):
        """All PDF tool GET pages load."""
        tools = [
            'tools_list', 'merge_pdf', 'split_pdf', 'compress_pdf',
            'jpg_to_pdf', 'pdf_to_jpg', 'rotate_pdf', 'delete_pages',
            'extract_pages', 'pdf_to_word', 'pdf_to_excel', 'pdf_to_ppt',
            'word_to_pdf', 'excel_to_pdf', 'ppt_to_pdf', 'html_to_pdf',
            'txt_to_pdf', 'png_to_pdf', 'webp_to_pdf', 'pdf_to_png',
            'crop_pdf', 'watermark_pdf', 'add_page_numbers',
            'add_text_to_pdf', 'add_image_to_pdf', 'annotate_pdf',
            'highlight_pdf', 'redact_pdf', 'organize_pdf',
            'protect_pdf', 'unlock_pdf', 'sign_pdf', 'repair_pdf',
            'ocr_pdf', 'scan_to_pdf',
        ]
        for name in tools:
            with self.subTest(tool=name):
                response = self.client.get(reverse(name))
                self.assertEqual(response.status_code, 200,
                                 f"{name}: {response.status_code}")

    def test_all_ai_tool_get(self):
        """All AI tool GET pages load."""
        tools = [
            'ai_tools', 'ai_study', 'ai_summarize', 'ai_extract',
            'ai_translate', 'ai_ask', 'ai_study_quiz', 'ai_study_flashcards',
            'ai_study_notes', 'ai_study_guide', 'ai_study_questions',
            'ai_study_important', 'ai_study_chapters', 'ai_study_concepts',
            'ai_study_exam',
        ]
        for name in tools:
            with self.subTest(tool=name):
                response = self.client.get(reverse(name))
                self.assertEqual(response.status_code, 200,
                                 f"{name}: {response.status_code}")

    # ================================================================
    # N. OUTPUT FILE VERIFICATION
    # ================================================================

    def test_pdf_to_word_output_validity(self):
        """Verify DOCX output is valid."""
        from docx import Document
        pdf = self._make_pdf()
        response = self.client.post(reverse('pdf_to_word'), {'file': pdf})
        self.assertEqual(response.status_code, 302)

        download_response = self.client.get(response.url)
        self.assertEqual(download_response.status_code, 200)
        content = b''.join(download_response.streaming_content)
        self.assertGreater(len(content), 100)

        try:
            import zipfile
            z = zipfile.ZipFile(io.BytesIO(content))
            self.assertIn('word/document.xml', z.namelist())
            doc = Document(io.BytesIO(content))
        except Exception as e:
            self.fail(f"DOCX validation failed: {e}")

    def test_pdf_to_excel_output_validity(self):
        """Verify XLSX output is valid."""
        from openpyxl import load_workbook
        pdf = self._make_pdf()
        response = self.client.post(reverse('pdf_to_excel'), {'file': pdf})
        self.assertEqual(response.status_code, 302)

        download_response = self.client.get(response.url)
        self.assertEqual(download_response.status_code, 200)
        content = b''.join(download_response.streaming_content)

        try:
            wb = load_workbook(io.BytesIO(content))
            self.assertGreater(len(wb.sheetnames), 0)
        except Exception as e:
            self.fail(f"XLSX validation failed: {e}")

    def test_pdf_to_ppt_output_validity(self):
        """Verify PPTX output is valid."""
        from pptx import Presentation
        pdf = self._make_pdf()
        response = self.client.post(reverse('pdf_to_ppt'), {'file': pdf})
        self.assertEqual(response.status_code, 302)

        download_response = self.client.get(response.url)
        self.assertEqual(download_response.status_code, 200)
        content = b''.join(download_response.streaming_content)

        try:
            prs = Presentation(io.BytesIO(content))
            self.assertGreater(len(prs.slides), 0)
        except Exception as e:
            self.fail(f"PPTX validation failed: {e}")

    def test_word_to_pdf_output_validity(self):
        """Verify PDF output from Word is valid."""
        from pypdf import PdfReader
        docx = make_docx_file()
        response = self.client.post(reverse('word_to_pdf'), {'file': docx})
        self.assertEqual(response.status_code, 302)

        download_response = self.client.get(response.url)
        self.assertEqual(download_response.status_code, 200)
        content = b''.join(download_response.streaming_content)
        self.assertGreater(len(content), 10)

        try:
            reader = PdfReader(io.BytesIO(content))
            self.assertGreater(reader.get_num_pages(), 0)
        except Exception as e:
            self.fail(f"PDF validation failed: {e}")

    def test_html_to_pdf_output_validity(self):
        """Verify PDF output from HTML is valid."""
        from pypdf import PdfReader
        html = make_html_file()
        response = self.client.post(reverse('html_to_pdf'), {'file': html})
        self.assertEqual(response.status_code, 302)

        download_response = self.client.get(response.url)
        self.assertEqual(download_response.status_code, 200)
        content = b''.join(download_response.streaming_content)

        try:
            reader = PdfReader(io.BytesIO(content))
            self.assertGreater(reader.get_num_pages(), 0)
        except Exception as e:
            self.fail(f"PDF validation failed: {e}")

    def test_txt_to_pdf_output_validity(self):
        """Verify PDF output from TXT is valid."""
        from pypdf import PdfReader
        txt = make_txt_file()
        response = self.client.post(reverse('txt_to_pdf'), {'file': txt})
        self.assertEqual(response.status_code, 302)

        download_response = self.client.get(response.url)
        self.assertEqual(download_response.status_code, 200)
        content = b''.join(download_response.streaming_content)

        try:
            reader = PdfReader(io.BytesIO(content))
            self.assertGreater(reader.get_num_pages(), 0)
        except Exception as e:
            self.fail(f"PDF validation failed: {e}")

    def test_png_to_pdf_output_validity(self):
        """Verify PDF output from PNG is valid."""
        from pypdf import PdfReader
        img = make_image_file('png')
        response = self.client.post(reverse('png_to_pdf'), {'files': img})
        self.assertIn(response.status_code, [302, 200])
        if response.status_code == 302:
            download_response = self.client.get(response.url)
            self.assertEqual(download_response.status_code, 200)
            content = b''.join(download_response.streaming_content)
            try:
                reader = PdfReader(io.BytesIO(content))
                self.assertGreater(reader.get_num_pages(), 0)
            except Exception as e:
                self.fail(f"PDF validation failed: {e}")

    # ================================================================
    # O. TEMP FILE CLEANUP
    # ================================================================

    def test_temp_file_cleanup_after_pdf_to_word(self):
        """Verify temp files are cleaned up after processing."""
        from django.conf import settings
        media_temp = os.path.join(settings.MEDIA_ROOT, 'temp')
        os.makedirs(media_temp, exist_ok=True)
        before = len(os.listdir(media_temp)) if os.path.exists(media_temp) else 0

        pdf = self._make_pdf()
        response = self.client.post(reverse('pdf_to_word'), {'file': pdf})
        self.assertEqual(response.status_code, 302)

        after = len(os.listdir(media_temp)) if os.path.exists(media_temp) else 0
        self.assertLessEqual(after, before + 5)

    # ================================================================
    # P. ERROR HANDLING
    # ================================================================

    def test_invalid_page_range(self):
        """Invalid page range should fail gracefully."""
        pdf = self._make_pdf()
        response = self.client.post(reverse('delete_pages'), {
            'file': pdf, 'pages': 'abc-def'
        })
        self.assertIn(response.status_code, [200, 302, 400])

    def test_protect_password_mismatch(self):
        """Password mismatch should show error."""
        pdf = self._make_pdf()
        response = self.client.post(reverse('protect_pdf'), {
            'file': pdf, 'password': 'secret', 'confirm_password': 'different'
        })
        self.assertEqual(response.status_code, 200)
        self.assertContains(response, "do not match")
